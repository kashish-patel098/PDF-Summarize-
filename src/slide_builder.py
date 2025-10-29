"""Create PPTX using python-pptx. Adds notes.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
import os


def build_pptx(slide_items, visuals, out_path: str):
    """slide_items: list of dicts with headline, bullets, note
       visuals: list of image paths corresponding to slides
    """
    os.makedirs(os.path.dirname(out_path) or '.', exist_ok=True)
    prs = Presentation()
    prs.slide_height = Inches(4.5)  # 720px @ 160 dpi roughly
    prs.slide_width = Inches(8)
    for i, item in enumerate(slide_items):
        layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(layout)
        # Use full slide width for content (no images). Title at top, bullets fill remaining space.
        # Title box (full width)
        left = Inches(0.3)
        top = Inches(0.2)
        width = prs.slide_width - Inches(0.6)
        height = Inches(1.2)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        # enable word wrap and tighten margins so long titles wrap inside the box
        try:
            tf.word_wrap = True
        except Exception:
            pass
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.06)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        # set headline text and font size using runs
        tf.text = item.get('headline', '')
        # ensure at least one paragraph exists
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                try:
                    run.font.size = Pt(32)
                except Exception:
                    pass

        # Content box uses the rest of the slide (full width)
        content_top = top + height + Inches(0.15)
        content_height = prs.slide_height - content_top - Inches(0.4)
        cb = slide.shapes.add_textbox(left, content_top, width, content_height)
        ctf = cb.text_frame
        try:
            ctf.word_wrap = True
        except Exception:
            pass
        # give the content more breathing room
        ctf.margin_left = Inches(0.06)
        ctf.margin_right = Inches(0.06)
        ctf.margin_top = Inches(0.03)
        ctf.margin_bottom = Inches(0.03)
        # ensure default paragraph is cleared
        if ctf.paragraphs and ctf.paragraphs[0].text:
            ctf.paragraphs[0].text = ''
        # add bullets (full-width, larger font for readability)
        for b in item.get('bullets', []):
            p = ctf.add_paragraph()
            p.text = b
            p.level = 0
            for run in p.runs:
                try:
                    run.font.size = Pt(18)
                except Exception:
                    pass
        # notes
        notes = slide.notes_slide.notes_text_frame
        notes.text = item.get('note', '')
    prs.save(out_path)
    return out_path

